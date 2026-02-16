import streamlit as st
import os
import re
import io
import json
import tempfile
import logging
from datetime import datetime
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import traceback
from PIL import Image, ImageEnhance, ImageFilter, ImageOps
import numpy as np

# Настройка логирования
logging.basicConfig(
    level=logging.WARNING,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('presentation_analyzer.log'),
    ]
)
logger = logging.getLogger(__name__)

# Проверка Tesseract
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
    try:
        langs = pytesseract.get_languages()
        if 'rus' in langs and 'eng' in langs:
            OCR_LANGUAGES = 'rus+eng'
        elif 'rus' in langs:
            OCR_LANGUAGES = 'rus'
        else:
            OCR_LANGUAGES = 'eng'
    except:
        OCR_LANGUAGES = 'rus+eng'
        
except Exception as e:
    TESSERACT_AVAILABLE = False
    OCR_LANGUAGES = 'rus+eng'

def deduplicate_pptx(pptx_path: str) -> None:
    """Пересобирает PPTX без дубликатов файлов внутри ZIP.
    PowerPoint иногда показывает 'ошибка при открытии' если в архиве есть повторяющиеся имена.
    """
    import zipfile, os, tempfile
    if not os.path.exists(pptx_path):
        return
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            # Берем последнюю версию каждого файла
            names = zin.namelist()
            last_index = {}
            for i, n in enumerate(names):
                last_index[n] = i
            keep = {n for n,i in last_index.items()}
            with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                for i, n in enumerate(names):
                    if last_index.get(n) != i:
                        continue
                    zout.writestr(n, zin.read(n))
        os.replace(tmp_path, pptx_path)
    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except:
            pass


class PresentationAnalyzer:
    def __init__(self, pptx_path):
        """Инициализация анализатора с путем к презентации"""
        self.pptx_path = pptx_path
        self.results = []
        self.used_fonts = set()
        self.ocr_languages = OCR_LANGUAGES
        self.full_ocr_texts = {}
        self.analysis_timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        self.selected_slides_range = 'all'
        
        # Настройки
        self.settings = {
            'text_on_image_threshold': 0.3,
            'min_text_length_for_ocr': 3,
            'ocr_confidence_threshold': 0.6,
            'overlap_threshold': 0.2,
            'max_text_chars': 1000,
            'min_image_area_percentage': 0.7,
            'ocr_min_confidence': 45,
            'ocr_alternate_min_confidence': 35,
            'max_ocr_text_length': 5000,
        }
    
    def analyze_selected_slides(self, slides_range='all'):
        """Анализ выбранных слайдов"""
        try:
            self.selected_slides_range = slides_range
            
            prs = Presentation(self.pptx_path)
            total_slides = len(prs.slides)
            
            # Определяем, какие слайды анализировать
            slides_to_analyze = self.parse_slides_range(slides_range, total_slides)
            
            if not slides_to_analyze:
                return [], {}
            
            presentation_stats = {
                'has_animations': False,
                'has_transitions': False,
                'fonts_count': 0,
                'background_issues': 0,
                'text_on_images': 0,
                'total_images': 0,
                'ocr_used': False,
                'ocr_text_found': 0,
                'total_ocr_characters': 0,
                'selected_slides_count': len(slides_to_analyze),
                'selected_slides_range': slides_range,
                'total_slides_in_presentation': total_slides,
            }
            
            # Проверяем переходы между слайдами
            has_transitions = self.check_presentation_transitions(prs)
            presentation_stats['has_transitions'] = has_transitions
            
            for slide_num in slides_to_analyze:
                i = slide_num
                slide_result = self.analyze_slide(prs.slides[i-1], i)
                self.results.append(slide_result)
                
                if slide_result['Анимации'] == '✗':
                    presentation_stats['has_animations'] = True
                if slide_result['Фон'] == '✗':
                    presentation_stats['background_issues'] += 1
                if slide_result['Текст_на_изобр'] == 'Да':
                    presentation_stats['text_on_images'] += 1
                    presentation_stats['ocr_text_found'] += 1
                if slide_result['Изображения'] > 0:
                    presentation_stats['total_images'] += slide_result['Изображения']
                
                if slide_result.get('OCR_текст'):
                    presentation_stats['total_ocr_characters'] += len(slide_result['OCR_текст'])
                    presentation_stats['ocr_used'] = True
            
            self.analyze_fonts()
            presentation_stats['fonts_count'] = len(self.used_fonts)
            
            # Сохраняем полные OCR результаты
            self.save_full_ocr_results()
            
            return self.results, presentation_stats
            
        except Exception as e:
            st.error(f"Ошибка анализа слайдов: {e}")
            traceback.print_exc()
            return [], {}
    
    def calculate_conformance_percentage(self, results, presentation_stats):
        """Расчет процента соответствия критериям"""
        try:
            total_slides = len(results)
            
            # Веса для каждого критерия
            weights = {
                'background': 15,      # Белый фон на всех слайдах
                'fonts': 15,           # Не более 2 шрифтов
                'text_overload': 10,   # Не более 1000 символов на слайде
                'text_on_images': 15,  # Нет текста на изображениях
                'animations': 15,      # Нет анимаций
                'transitions': 10,     # Нет переходов
                'slide_compliance': 20 # Общее соответствие каждого слайда
            }
            
            total_possible = sum(weights.values())
            achieved_score = 0
            
            # 1. Критерий фона
            background_issues = presentation_stats.get('background_issues', 0)
            if total_slides > 0:
                background_score = (total_slides - background_issues) / total_slides * weights['background']
            else:
                background_score = weights['background']
            achieved_score += background_score
            
            # 2. Критерий шрифтов
            fonts_count = presentation_stats.get('fonts_count', 0)
            if fonts_count <= 2:
                fonts_score = weights['fonts']
            elif fonts_count <= 3:
                fonts_score = weights['fonts'] * 0.5
            else:
                fonts_score = 0
            achieved_score += fonts_score
            
            # 3. Критерий текстовой перегрузки
            text_issues = sum(1 for r in results if r['Текст'] == '✗')
            if total_slides > 0:
                text_score = (total_slides - text_issues) / total_slides * weights['text_overload']
            else:
                text_score = weights['text_overload']
            achieved_score += text_score
            
            # 4. Критерий текста на изображениях
            text_on_images = presentation_stats.get('text_on_images', 0)
            if total_slides > 0:
                images_score = (total_slides - text_on_images) / total_slides * weights['text_on_images']
            else:
                images_score = weights['text_on_images']
            achieved_score += images_score
            
            # 5. Критерий анимаций
            anim_issues = sum(1 for r in results if r['Анимации'] == '✗')
            if total_slides > 0:
                anim_score = (total_slides - anim_issues) / total_slides * weights['animations']
            else:
                anim_score = weights['animations']
            achieved_score += anim_score
            
            # 6. Критерий переходов
            transition_issues = 1 if presentation_stats.get('has_transitions') else 0
            if transition_issues == 0:
                transition_score = weights['transitions']
            else:
                transition_score = 0
            achieved_score += transition_score
            
            # 7. Критерий соответствия каждого слайда
            compliant_slides = 0
            for result in results:
                is_compliant = (
                    result['Фон'] == '✓' and
                    result['Шрифты'] == '✓' and
                    result['Текст'] == '✓' and
                    result['Текст_на_изобр'] == 'Нет' and
                    result['Анимации'] == '✓'
                )
                if is_compliant:
                    compliant_slides += 1
            
            if total_slides > 0:
                slide_score = (compliant_slides / total_slides) * weights['slide_compliance']
            else:
                slide_score = weights['slide_compliance']
            achieved_score += slide_score
            
            # Расчет процента
            percentage = (achieved_score / total_possible) * 100
            percentage = round(percentage, 1)
            
            # Определяем уровень готовности
            if percentage >= 90:
                readiness_level = "отлично"
                readiness_color = "#27ae60"
                readiness_emoji = "🎉"
            elif percentage >= 75:
                readiness_level = "хорошо"
                readiness_color = "#2ecc71"
                readiness_emoji = "👍"
            elif percentage >= 60:
                readiness_level = "удовлетворительно"
                readiness_color = "#f39c12"
                readiness_emoji = "⚠️"
            elif percentage >= 40:
                readiness_level = "требует доработки"
                readiness_color = "#e74c3c"
                readiness_emoji = "🔧"
            else:
                readiness_level = "критически низкая"
                readiness_color = "#c0392b"
                readiness_emoji = "🚨"
            
            # Определяем можно ли отправлять дизайнерам
            can_send_to_designers = percentage >= 57
            
            # Детальная информация по критериям
            criteria_details = {
                'background': {
                    'score': round(background_score, 1),
                    'max': weights['background'],
                    'issues': background_issues,
                    'total': total_slides,
                    'description': 'Белый фон на всех слайдах'
                },
                'fonts': {
                    'score': round(fonts_score, 1),
                    'max': weights['fonts'],
                    'fonts_count': fonts_count,
                    'description': 'Не более 2 шрифтов в презентации'
                },
                'text_overload': {
                    'score': round(text_score, 1),
                    'max': weights['text_overload'],
                    'issues': text_issues,
                    'total': total_slides,
                    'description': 'Не более 1000 символов на слайде'
                },
                'text_on_images': {
                    'score': round(images_score, 1),
                    'max': weights['text_on_images'],
                    'issues': text_on_images,
                    'total': total_slides,
                    'description': 'Нет текста на изображениях'
                },
                'animations': {
                    'score': round(anim_score, 1),
                    'max': weights['animations'],
                    'issues': anim_issues,
                    'total': total_slides,
                    'description': 'Нет анимаций на слайдах'
                },
                'transitions': {
                    'score': round(transition_score, 1),
                    'max': weights['transitions'],
                    'has_issues': transition_issues > 0,
                    'description': 'Нет переходов между слайдами'
                },
                'slide_compliance': {
                    'score': round(slide_score, 1),
                    'max': weights['slide_compliance'],
                    'compliant': compliant_slides,
                    'total': total_slides,
                    'description': 'Полностью соответствующие слайды'
                }
            }
            
            # Генерация рекомендаций
            recommendations = []
            if percentage < 57:
                recommendations.append("Рекомендуется доработать презентацию перед отправкой дизайнерам")
            if background_issues > 0:
                recommendations.append(f"Исправьте фон на {background_issues} слайдах")
            if fonts_count > 2:
                recommendations.append(f"Уменьшите количество шрифтов с {fonts_count} до 2")
            if text_issues > 0:
                recommendations.append(f"Уменьшите текст на {text_issues} слайдах")
            if text_on_images > 0:
                recommendations.append(f"Уберите текст с изображений на {text_on_images} слайдах")
            if anim_issues > 0:
                recommendations.append(f"Удалите анимации с {anim_issues} слайдов")
            if transition_issues > 0:
                recommendations.append("Удалите переходы между слайдами")
            
            # Текст вывода для пользователя
            if can_send_to_designers:
                user_message = f"🎉 Ваша презентация соответствует критериям на {percentage}%. Презентация готова для отправки дизайнерам!"
            else:
                user_message = f"⚠️ Ваша презентация соответствует критериям на {percentage}%. Если Вы планируете отправлять дизайнерам, рекомендуется её доработать."
            
            conformance_info = {
                'percentage': percentage,
                'readiness_level': readiness_level,
                'readiness_color': readiness_color,
                'readiness_emoji': readiness_emoji,
                'can_send_to_designers': can_send_to_designers,
                'criteria_details': criteria_details,
                'recommendations': recommendations,
                'user_message': user_message,
                'total_possible_score': total_possible,
                'achieved_score': round(achieved_score, 1),
                'compliant_slides': compliant_slides,
                'total_slides': total_slides
            }
            
            return conformance_info
            
        except Exception as e:
            st.error(f"Ошибка расчета процента соответствия: {e}")
            return None
    
    def parse_slides_range(self, slides_range, total_slides):
        """Парсинг диапазона слайдов с поддержкой сложных форматов"""
        slides_to_analyze = []
        
        try:
            # Если значение пустое или 'all', возвращаем все слайды
            if not slides_range or str(slides_range).lower() == 'all':
                return list(range(1, total_slides + 1))
            
            slides_range = str(slides_range).strip()
            
            # Если это одно число
            if slides_range.isdigit():
                slide_num = int(slides_range)
                if 1 <= slide_num <= total_slides:
                    return [slide_num]
                else:
                    return []
            
            # Удаляем все пробелы для упрощения обработки
            slides_range = slides_range.replace(' ', '')
            
            # Проверяем наличие запятых (список слайдов)
            if ',' in slides_range:
                parts = slides_range.split(',')
                for part in parts:
                    if '-' in part:
                        # Диапазон внутри списка (например: "1-3,5-7")
                        range_parts = part.split('-')
                        if len(range_parts) == 2 and range_parts[0].isdigit() and range_parts[1].isdigit():
                            start = int(range_parts[0])
                            end = int(range_parts[1])
                            slides_to_analyze.extend(range(start, min(end, total_slides) + 1))
                    elif part.isdigit():
                        # Один номер слайда
                        slide_num = int(part)
                        if 1 <= slide_num <= total_slides:
                            slides_to_analyze.append(slide_num)
            
            # Проверяем наличие дефиса (диапазон слайдов)
            elif '-' in slides_range:
                parts = slides_range.split('-')
                if len(parts) == 2 and parts[0].isdigit() and parts[1].isdigit():
                    start = int(parts[0])
                    end = int(parts[1])
                    slides_to_analyze = list(range(start, min(end, total_slides) + 1))
            
            # Удаляем дубликаты и сортируем
            slides_to_analyze = sorted(set(slides_to_analyze))
            
            # Если после всех проверок список пуст, возвращаем все слайды
            if not slides_to_analyze:
                slides_to_analyze = list(range(1, total_slides + 1))
                self.selected_slides_range = 'all'
            
        except Exception as e:
            st.error(f"Ошибка парсинга диапазона слайдов: {e}")
            slides_to_analyze = list(range(1, total_slides + 1))
            self.selected_slides_range = 'all'
        
        return slides_to_analyze
    
    def generate_word_report(self, results, presentation_stats, output_path=None):
        """Генерация отчета в формате Word (.docx)"""
        try:
            from docx import Document
            from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # Заголовок
            title = doc.add_heading('Отчет анализа презентации', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Информация о файле
            doc.add_paragraph(f'Файл: {os.path.basename(self.pptx_path)}')
            doc.add_paragraph(f'Дата анализа: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
            doc.add_paragraph(f'Всего слайдов в презентации: {presentation_stats.get("total_slides_in_presentation", len(results))}')
            doc.add_paragraph(f'Проанализировано слайдов: {len(results)}')
            doc.add_paragraph(f'Диапазон анализа: {self.selected_slides_range}')
            
            # Рассчитываем процент соответствия
            conformance_info = self.calculate_conformance_percentage(results, presentation_stats)
            
            if conformance_info:
                doc.add_paragraph()
                doc.add_heading('Уровень соответствия критериям', level=1)
                
                # Создаем таблицу для процента соответствия
                conformance_table = doc.add_table(rows=2, cols=2)
                conformance_table.style = 'LightShading-Accent1'
                
                # Заголовок
                header_cells = conformance_table.rows[0].cells
                header_cells[0].text = 'Показатель'
                header_cells[1].text = 'Значение'
                
                # Данные
                row_cells = conformance_table.rows[1].cells
                row_cells[0].text = 'Соответствие критериям'
                row_cells[1].text = f"{conformance_info['percentage']}% ({conformance_info['readiness_level']})"
                
                # Добавляем рекомендацию
                doc.add_paragraph()
                if conformance_info['can_send_to_designers']:
                    doc.add_paragraph(f"✅ {conformance_info['user_message']}")
                else:
                    doc.add_paragraph(f"⚠️ {conformance_info['user_message']}")
                
                # Добавляем рекомендации
                if conformance_info['recommendations']:
                    doc.add_paragraph()
                    doc.add_heading('Рекомендации по улучшению:', level=2)
                    for rec in conformance_info['recommendations']:
                        para = doc.add_paragraph()
                        para.add_run('• ').bold = False
                        para.add_run(rec)
            
            doc.add_paragraph()
            
            # 1. РЕКОМЕНДАЦИИ
            doc.add_heading('1. Рекомендации', level=1)
            
            recommendations = [
                "Используйте не более 2 шрифтов в презентации",
                "Убедитесь, что фон всех слайдов соответствует шаблону",
                "Избегайте анимаций и переходов",
                "Не размещайте текст на изображениях",
                "Ограничьте текст на слайде 1000 символами",
                "Используйте редактируемые текстовые блоки вместо текста в изображениях"
            ]
            
            for rec in recommendations:
                paragraph = doc.add_paragraph()
                paragraph.add_run('• ').bold = False
                paragraph.add_run(rec)
            
            doc.add_paragraph()
            
            # 2. СВОДНАЯ СТАТИСТИКА
            doc.add_heading('2. Сводная статистика', level=1)
            
            stats_data = [
                ('Всего слайдов', len(results)),
                ('Слайдов с нарушениями', sum(1 for r in results if r['Статус'] != 'OK')),
                ('Фон не соответствует шаблону', presentation_stats.get('background_issues', 0)),
                ('Использовано шрифтов', presentation_stats.get('fonts_count', 0)),
                ('Текст на изображениях', presentation_stats.get('text_on_images', 0)),
                ('Всего изображений', presentation_stats.get('total_images', 0)),
                ('Анимации', 'Да' if presentation_stats.get('has_animations') else 'Нет'),
                ('Переходы', 'Да' if presentation_stats.get('has_transitions') else 'Нет'),
            ]
            
            stats_table = doc.add_table(rows=len(stats_data), cols=2)
            stats_table.style = 'LightShading-Accent1'
            
            for i, (key, value) in enumerate(stats_data):
                row_cells = stats_table.rows[i].cells
                row_cells[0].text = str(key)
                row_cells[1].text = str(value)
            
            doc.add_paragraph()
            doc.add_heading('3. Детальный анализ по слайдам', level=1)
            
            # Таблица для детального анализа
            slides_table = doc.add_table(rows=1, cols=9)
            slides_table.style = 'Table Grid'
            
            headers = ['Слайд', 'Наличие отклонения', 'Фон', 'Шрифты', 'Текст', 'Элементы', 'Изображения', 'Текст на изобр.', 'Анимации']
            header_cells = slides_table.rows[0].cells
            
            for i, header in enumerate(headers):
                header_cells[i].text = header
                header_cells[i].paragraphs[0].runs[0].bold = True
                header_cells[i].paragraphs[0].runs[0].font.size = DocxPt(9)
            
            # Заполняем данные
            for result in results:
                row_cells = slides_table.add_row().cells
                
                # Слайд номер
                row_cells[0].text = str(result['Слайд'])
                row_cells[0].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Наличие отклонения
                has_violation = result['Статус'] != 'OK'
                row_cells[1].text = 'Да' if has_violation else 'Нет'
                if has_violation:
                    row_cells[1].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[1].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Фон
                row_cells[2].text = result['Фон']
                if result['Фон'] == '✗':
                    row_cells[2].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[2].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Шрифты
                row_cells[3].text = result['Шрифты']
                if result['Шрифты'] == '✗':
                    row_cells[3].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[3].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Текст
                row_cells[4].text = result['Текст_дет']
                if result['Текст'] == '✗':
                    row_cells[4].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[4].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Элементы
                row_cells[5].text = str(result['Элементы'])
                row_cells[5].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Изображения
                row_cells[6].text = str(result['Изображения'])
                row_cells[6].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Текст на изображениях
                row_cells[7].text = result['Текст_на_изобр']
                if result['Текст_на_изобр'] == 'Да':
                    row_cells[7].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[7].paragraphs[0].runs[0].font.size = DocxPt(9)
                
                # Анимации
                row_cells[8].text = result['Анимации']
                if result['Анимации'] == '✗':
                    row_cells[8].paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 0, 0)
                row_cells[8].paragraphs[0].runs[0].font.size = DocxPt(9)
            
            # 4. Текст, найденный на изображениях (OCR)
            if self.full_ocr_texts:
                doc.add_paragraph()
                doc.add_heading('4. Текст, найденный на изображениях (OCR)', level=1)
                
                for slide_num, ocr_data in self.full_ocr_texts.items():
                    doc.add_heading(f'Слайд {slide_num}', level=2)
                    doc.add_paragraph(f'Изображений на слайде: {ocr_data.get("image_count", 0)}')
                    doc.add_paragraph(f'Изображений с текстом: {ocr_data.get("images_with_text", 0)}')
                    doc.add_paragraph(f'Уверенность: {ocr_data.get("confidence", 0):.1f}%')
                    doc.add_paragraph(f'Метод: {ocr_data.get("method", "unknown")}')
                    
                    ocr_text = ocr_data.get('text', '')
                    if ocr_text:
                        if len(ocr_text) > 5000:
                            ocr_text = ocr_text[:5000] + '\n... [текст сокращен]'
                        
                        formatted_text = ocr_text.replace('--- Текст с изображения', '\n--- Текст с изображения')
                        text_paragraph = doc.add_paragraph(formatted_text)
                        text_paragraph.style = 'Normal'
            
            # Сохраняем документ
            if output_path is None:
                output_path = f"report_{self.analysis_timestamp}.docx"
            
            doc.save(output_path)
            return output_path
            
        except Exception as e:
            st.error(f"Ошибка генерации Word отчета: {e}")
            traceback.print_exc()
            return None
    
    def save_full_ocr_results(self):
        """Сохранение полных OCR результатов"""
        try:
            for slide_result in self.results:
                if slide_result.get('OCR_текст'):
                    slide_num = slide_result['Слайд']
                    self.full_ocr_texts[slide_num] = {
                        'text': slide_result.get('OCR_текст', ''),
                        'confidence': slide_result.get('OCR_уверенность', 0),
                        'method': slide_result.get('OCR_метод', ''),
                        'image_count': slide_result['Изображения'],
                        'images_with_text': slide_result.get('OCR_изображений_с_текстом', 0)
                    }
            
        except Exception as e:
            pass
    
    def check_presentation_transitions(self, prs):
        """Проверка переходов между слайдами"""
        try:
            for slide in prs.slides:
                try:
                    slide_xml = str(slide.element.xml).lower()
                    if any(keyword in slide_xml for keyword in ['p:transition', 'transition']):
                        return True
                except:
                    continue
        except:
            pass
        return False
    
    def analyze_slide(self, slide, slide_num):
        """Анализ одного слайда"""
        slide_result = {
            'Слайд': slide_num,
            'Статус': 'OK',
            'Нарушения': [],
            'Шрифты': '✓',
            'Текст': '✓',
            'Анимации': '✓',
            'Переходы': '✓',
            'Фон': '✓',
            'Изображения': 0,
            'Текст_на_изобр': 'Нет',
            'Текст_дет': '',
            'Элементы': len(slide.shapes),
            'OCR_текст': '',
            'OCR_уверенность': 0,
            'OCR_метод': '',
            'OCR_изображений_с_текстом': 0,
        }
        
        # 1. Проверка фона
        if not self.check_background_comprehensive(slide):
            slide_result['Фон'] = '✗'
            slide_result['Нарушения'].append('ФОН')
        
        # 2. Проверка текста
        text_overload, char_count = self.check_text_improved(slide)
        slide_result['Текст_дет'] = f"{char_count} симв."
        if text_overload:
            slide_result['Текст'] = '✗'
            slide_result['Нарушения'].append(f'ТЕКСТ({char_count})')
        
        # 3. Проверка анимаций
        if self.check_animations_improved(slide):
            slide_result['Анимации'] = '✗'
            slide_result['Нарушения'].append('АНИМАЦИИ')
        
        # 4. Проверка изображений
        has_text_on_images, image_count, ocr_data = self.check_images_enhanced(slide)
        slide_result['Изображения'] = image_count
        
        if ocr_data:
            full_text = ocr_data.get('text', '')
            slide_result['OCR_текст'] = full_text[:self.settings['max_ocr_text_length']]
            slide_result['OCR_уверенность'] = ocr_data.get('confidence', 0)
            slide_result['OCR_метод'] = ocr_data.get('method', '')
            slide_result['OCR_изображений_с_текстом'] = ocr_data.get('images_with_text', 0)
        
        if has_text_on_images:
            slide_result['Текст_на_изобр'] = 'Да'
            slide_result['Нарушения'].append('ТЕКСТ_НА_ИЗОБР')
        
        # 5. Сбор шрифтов
        self.collect_fonts(slide)
        
        if slide_result['Нарушения']:
            slide_result['Статус'] = ', '.join(slide_result['Нарушения'])
        
        return slide_result
    
    def check_background_comprehensive(self, slide):
        """Комплексная проверка фона"""
        try:
            # Проверка фона слайда
            if slide.background:
                fill = slide.background.fill
                if fill.type == 1:
                    if hasattr(fill.fore_color, 'rgb'):
                        color = fill.fore_color.rgb
                        if hasattr(color, 'r'):
                            if not (color.r == 255 and color.g == 255 and color.b == 255):
                                return False
                        elif color != RGBColor(255, 255, 255):
                            return False
                elif fill.type != 0:
                    return False
            
            # Проверка крупных фигур
            try:
                slide_width = slide.width if hasattr(slide, 'width') else Inches(10)
                slide_height = slide.height if hasattr(slide, 'height') else Inches(7.5)
                slide_area = slide_width * slide_height
                
                for shape in slide.shapes:
                    try:
                        shape_area = shape.width * shape.height
                        if shape_area > slide_area * self.settings['min_image_area_percentage']:
                            if hasattr(shape, 'fill'):
                                fill = shape.fill
                                if fill.type == 1:
                                    if hasattr(fill.fore_color, 'rgb'):
                                        color = fill.fore_color.rgb
                                        if hasattr(color, 'r'):
                                            if not (color.r == 255 and color.g == 255 and color.b == 255):
                                                return False
                                        elif color != RGBColor(255, 255, 255):
                                            return False
                    except:
                        continue
            except:
                pass
            
            # Проверка XML на цвета
            try:
                slide_xml = str(slide.element.xml).lower()
                
                hex_pattern = r'#[0-9a-f]{6}'
                hex_matches = re.findall(hex_pattern, slide_xml)
                for hex_color in hex_matches:
                    if hex_color != '#ffffff' and hex_color != '#ffffff00':
                        return False
                
                rgb_pattern = r'rgb\((\d+),\s*(\d+),\s*(\d+)\)'
                rgb_matches = re.findall(rgb_pattern, slide_xml)
                for r, g, b in rgb_matches:
                    if int(r) != 255 or int(g) != 255 or int(b) != 255:
                        return False
            except:
                pass
            
            return True
            
        except:
            return False
    
    def check_text_improved(self, slide):
        """Проверка текста"""
        try:
            total_chars = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    text = shape.text_frame.text.strip()
                    if text and len(text) > 1:
                        clean_text = re.sub(r'\s+', ' ', text)
                        total_chars += len(clean_text)
            
            return total_chars > self.settings['max_text_chars'], total_chars
        except:
            return False, 0
    
    def check_animations_improved(self, slide):
        """Проверка анимаций"""
        try:
            xml = str(slide.element.xml).lower()
            animation_patterns = [
                r'<p:anim\s', r'p:ctn', r'p:seq', r'p:par',
                r'dur=["\']', r'accel=["\']', r'decel=["\']',
                r'<p:custanim\s', r'<p:set\s', r'animate\s',
                r'animation\s', r'animbullet\s', r'animeffect\s'
            ]
            
            for pattern in animation_patterns:
                if re.search(pattern, xml):
                    return True
            
            return False
        except:
            return False
    
    def check_images_enhanced(self, slide):
        """Проверка изображений и текста на них"""
        try:
            image_count = 0
            has_text_on_images = False
            ocr_data = {}
            
            image_info = []
            text_shapes = []
            
            def process_shape(shape):
                nonlocal image_count
                
                if hasattr(shape, 'shapes'):
                    for subshape in shape.shapes:
                        process_shape(subshape)
                    return
                
                if hasattr(shape, "image"):
                    image_count += 1
                    try:
                        img_info = {
                            'shape': shape,
                            'id': id(shape),
                            'index': image_count,
                            'width': shape.width,
                            'height': shape.height,
                            'format': shape.image.ext,
                        }
                        image_info.append(img_info)
                    except:
                        return
                
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if text and text.strip():
                        try:
                            text_shape_info = {
                                'shape': shape,
                                'id': id(shape),
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'right': shape.left + shape.width,
                                'bottom': shape.top + shape.height,
                                'text': text.strip(),
                                'char_count': len(text.strip())
                            }
                            text_shapes.append(text_shape_info)
                        except:
                            return
            
            for shape in slide.shapes:
                process_shape(shape)
            
            if not image_info:
                return False, 0, None
            
            # Геометрическая проверка
            for text_shape in text_shapes:
                if text_shape['char_count'] < self.settings['min_text_length_for_ocr']:
                    continue
                    
                for img in image_info:
                    if self.shapes_overlap_improved(text_shape, img):
                        has_text_on_images = True
                        break
                
                if has_text_on_images:
                    break
            
            # OCR проверка
            if TESSERACT_AVAILABLE and len(image_info) > 0:
                try:
                    ocr_results = self.check_images_with_multiple_ocr_methods(slide, image_info)
                    
                    combined_text = ""
                    total_confidence = 0
                    images_with_text = 0
                    best_method = ""
                    best_confidence = 0
                    
                    img_info_dict = {img['id']: img for img in image_info}
                    
                    for img_id, (ocr_text, confidence, img_format, method) in ocr_results.items():
                        if self.is_meaningful_text(ocr_text) and confidence > self.settings['ocr_alternate_min_confidence']:
                            if combined_text:
                                combined_text += f"\n\n--- Текст с изображения {img_info_dict[img_id]['index']} ({img_info_dict[img_id]['width']:.0f}x{img_info_dict[img_id]['height']:.0f}) ---\n"
                            else:
                                combined_text += f"--- Текст с изображения {img_info_dict[img_id]['index']} ({img_info_dict[img_id]['width']:.0f}x{img_info_dict[img_id]['height']:.0f}) ---\n"
                            
                            combined_text += ocr_text
                            total_confidence += confidence
                            images_with_text += 1
                            
                            if confidence > best_confidence:
                                best_confidence = confidence
                                best_method = method
                    
                    if combined_text:
                        has_text_on_images = True
                        avg_confidence = total_confidence / images_with_text if images_with_text > 0 else 0
                        
                        ocr_data = {
                            'text': combined_text,
                            'confidence': avg_confidence,
                            'method': best_method if best_method else "multiple",
                            'image_count': len(image_info),
                            'images_with_text': images_with_text
                        }
                        
                except Exception as e:
                    pass
            
            return has_text_on_images, len(image_info), ocr_data
            
        except Exception as e:
            return False, 0, None
    
    def shapes_overlap_improved(self, shape1, shape2):
        """Проверка пересечения фигур"""
        try:
            overlap_x = not (shape1['right'] <= shape2['left'] or shape1['left'] >= shape2['right'])
            overlap_y = not (shape1['bottom'] <= shape2['top'] or shape1['top'] >= shape2['bottom'])
            return overlap_x and overlap_y
        except:
            return False
    
    def check_images_with_multiple_ocr_methods(self, slide, image_info):
        """OCR проверка с несколькими методами"""
        try:
            ocr_results = {}
            
            with tempfile.TemporaryDirectory() as temp_dir:
                for i, img_info in enumerate(image_info):
                    try:
                        shape = img_info['shape']
                        
                        if shape.width < 50 or shape.height < 50:
                            continue
                        
                        image_data = shape.image.blob
                        img_format = img_info.get('format', 'unknown')
                        
                        best_result = self.try_multiple_ocr_methods(image_data, img_format, img_info['index'], img_info['id'])
                        
                        if best_result:
                            text, confidence, method_used = best_result
                            ocr_results[img_info['id']] = (text, confidence, img_format, method_used)
                            
                    except Exception as e:
                        continue
            
            return ocr_results
            
        except Exception as e:
            return {}
    
    def try_multiple_ocr_methods(self, image_data, img_format, index, img_id):
        """Пробуем несколько методов OCR"""
        best_text = ""
        best_confidence = 0
        best_method = ""
        
        ocr_methods = [
            {'name': 'Tesseract_PSM6_rus+eng', 'config': '--oem 3 --psm 6 -l rus+eng', 'preprocess': 'standard'},
            {'name': 'Tesseract_PSM3_rus+eng', 'config': '--oem 3 --psm 3 -l rus+eng', 'preprocess': 'standard'},
            {'name': 'Tesseract_PSM11_rus+eng', 'config': '--oem 3 --psm 11 -l rus+eng', 'preprocess': 'high_contrast'},
        ]
        
        for method in ocr_methods:
            try:
                processed_image = self.preprocess_for_ocr_method(image_data, img_format, index, img_id, method['preprocess'])
                
                if processed_image is None:
                    continue
                
                data = pytesseract.image_to_data(processed_image, config=method['config'], output_type=pytesseract.Output.DICT)
                
                text_parts = []
                confidences = []
                
                for j in range(len(data['text'])):
                    text_item = data['text'][j].strip()
                    if text_item and len(text_item) > 1:
                        text_parts.append(text_item)
                        if data['conf'][j] != '-1':
                            confidences.append(float(data['conf'][j]))
                
                if text_parts:
                    text = ' '.join(text_parts).strip()
                    avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                    
                    text = self.clean_ocr_text(text)
                    
                    if text and avg_confidence > best_confidence:
                        if self.quick_text_quality_check(text, avg_confidence):
                            best_text = text
                            best_confidence = avg_confidence
                            best_method = method['name']
                            
            except Exception as e:
                continue
        
        if best_text and best_confidence > self.settings['ocr_alternate_min_confidence']:
            return best_text, best_confidence, best_method
        
        return None
    
    def preprocess_for_ocr_method(self, image_data, img_format, index, img_id, method='standard'):
        """Предобработка изображения"""
        try:
            img = Image.open(io.BytesIO(image_data))
            
            if img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[3])
                else:
                    background.paste(img)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            img = img.convert('L')
            
            if method == 'standard':
                enhancer = ImageEnhance.Sharpness(img)
                img = enhancer.enhance(2.0)
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(1.5)
                img = ImageOps.autocontrast(img, cutoff=2)
            elif method == 'high_contrast':
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(3.0)
                img = ImageOps.autocontrast(img, cutoff=5)
                threshold = 200
                img = img.point(lambda p: 255 if p > threshold else 0)
            elif method == 'inverted':
                img = ImageOps.invert(img)
                enhancer = ImageEnhance.Sharpness(img)
                img = enhancer.enhance(2.0)
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(1.5)
                img = ImageOps.autocontrast(img, cutoff=2)
            
            return img
            
        except Exception as e:
            return None
    
    def clean_ocr_text(self, text):
        """Очистка OCR текста"""
        if not text:
            return text
        
        lines = text.split('\n')
        clean_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            alpha_count = sum(1 for c in line if c.isalpha())
            total_chars = len(line)
            
            if total_chars > 0 and alpha_count / total_chars > 0.3:
                clean_lines.append(line)
        
        text = '\n'.join(clean_lines)
        
        replacements = {
            'Сберё': 'Сбер', 'СберЪ': 'Сбер', 'сберё': 'сбер',
            'СБЕРё': 'СБЕР', 'ё': 'е', 'Ё': 'Е',
            '""': '"', "''": "'", '``': '`', '``': '"',
            '”': '"', '„': '"', '«': '"', '»': '"',
            '—': '-', '–': '-', '`': "'", '´': "'",
            '‘': "'", '’': "'",
        }
        
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        text = re.sub(r'\s[а-яА-Яa-zA-Z]\s', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    def quick_text_quality_check(self, text, confidence):
        """Проверка качества текста"""
        if not text or len(text) < 10:
            return False
        
        russian_letters = sum(1 for c in text if 'а' <= c.lower() <= 'я' or c in 'ёе')
        total_letters = sum(1 for c in text if c.isalpha())
        
        if total_letters == 0:
            return False
        
        russian_ratio = russian_letters / total_letters if total_letters > 0 else 0
        
        if confidence < 50 and russian_ratio < 0.8:
            return False
        elif russian_ratio < 0.5:
            return False
        
        russian_words = re.findall(r'\b[А-Яа-яёЁ]{3,}\b', text)
        if len(russian_words) < 2:
            return False
        
        for word in russian_words:
            if len(word) > 10:
                if any(word.count(word[i:i+3]) > 2 for i in range(len(word)-2)):
                    return False
        
        return True
    
    def is_meaningful_text(self, text):
        """Проверка на значимость текста"""
        if not text:
            return False
        
        text = self.clean_ocr_text(text)
        
        if len(text) < 20:
            return False
        
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        if len(lines) < 1:
            return False
        
        meaningful_lines = 0
        
        for line in lines:
            if len(line) < 5:
                continue
            
            russian_letters = sum(1 for c in line if 'а' <= c.lower() <= 'я' or c in 'ёе')
            total_chars = len(line)
            russian_words = re.findall(r'\b[А-Яа-яёЁ]{3,}\b', line)
            
            is_text_like = (
                russian_letters > 5 and
                len(russian_words) > 1 and
                russian_letters / total_chars > 0.4
            )
            
            if is_text_like:
                meaningful_lines += 1
        
        return meaningful_lines >= 1
    
    def collect_fonts(self, slide):
        """Сбор шрифтов"""
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run.font, 'name') and run.font.name:
                                font_name = run.font.name
                                if font_name and font_name.strip():
                                    self.used_fonts.add(font_name.strip())
        except:
            pass
    
    def analyze_fonts(self):
        """Анализ шрифтов"""
        try:
            filtered_fonts = set()
            system_fonts = [
                '+mj-lt', '+mn-lt', 'calibri', 'tahoma', 'arial', 
                'times', 'verdana', 'cambria', 'segoe ui', 'consolas',
                'courier new', 'georgia', 'impact', 'trebuchet ms'
            ]
            
            for font in self.used_fonts:
                font_lower = font.lower()
                is_system_font = False
                
                for sys_font in system_fonts:
                    if sys_font in font_lower:
                        is_system_font = True
                        break
                
                if not is_system_font:
                    filtered_fonts.add(font)
            
            font_count = len(filtered_fonts)
            
            for result in self.results:
                if font_count > 2:
                    result['Шрифты'] = '✗'
                    if 'ШРИФТЫ' not in result['Нарушения']:
                        result['Нарушения'].append(f'ШРИФТЫ({font_count})')
                        result['Статус'] = ', '.join(result['Нарушения'])
            
        except:
            pass