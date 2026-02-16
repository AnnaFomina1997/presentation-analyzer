import streamlit as st
import os
import re
import io
import tempfile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

class PresentationGenerator:
    """Генератор исправленной презентации.

    Требования:
    - Титульники 1:1 как в template: используем слайды 1–3 из template и ТОЛЬКО подставляем текст.
    - Начиная с 4-го слайда используем фон/макет как у слайда 4 template (белый фон).
    - Для обычного текста: шрифт Montserrat, но РАЗМЕР берём из исходной презентации (чтобы ничего не "уезжало").
    - Изображения переносим с сохранением crop/rotate/effects через перенос XML + корректные rels (без "поломки rId").
    - Таблицы переносим с сохранением размеров шрифта из исходника (но шрифт Montserrat).
    """

    def __init__(self, pptx_path: str, template_path: str):
        self.pptx_path = pptx_path
        self.template_path = template_path

    # -----------------------------
    # Shape filtering (to avoid invisible "junk" shapes)
    # -----------------------------
    def _shape_has_meaningful_text(self, shape) -> bool:
        try:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                return bool(shape.text_frame.text and shape.text_frame.text.strip())
        except Exception:
            return False
        return False

    def _shape_has_visible_fill_or_line(self, shape):
        """Более 'добрая' проверка видимости, чтобы не терять важные фигуры.

        В презентациях часто встречаются заливки/линии из темы (fill.type == None),
        градиенты, schemeColor и т.п. Старые эвристики считали такие фигуры невидимыми,
        из‑за чего пропадали цветные прямоугольники и рамки.
        """
        try:
            # если это плейсхолдер и он пустой — почти всегда мусор
            if getattr(shape, "is_placeholder", False):
                if getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip():
                    return True
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    return True
                if getattr(shape, "has_table", False):
                    return True
                return False
        except Exception:
            pass

        # Если не удаётся корректно определить — лучше СКОПИРОВАТЬ, чем потерять
        try:
            # Заливка
            try:
                fill = shape.fill
                # fill.type == None часто означает "тематическая" заливка → видимо
                if getattr(fill, "type", None) is None:
                    return True
                if getattr(fill, "type", None) is not None:
                    # SOLID / GRADIENT / PATTERN и т.п.
                    return True
            except Exception:
                return True

            # Линия
            try:
                line = shape.line
                if line is not None:
                    w = getattr(line, "width", None)
                    if w is None:
                        return True
                    try:
                        return float(w) > 0
                    except Exception:
                        return True
            except Exception:
                return True

            return False
        except Exception:
            return True
    
    def _is_shape_meaningful(self, shape) -> bool:
        """Решает, копировать ли фигуру.

        Цель — НЕ терять важные фигуры. Поэтому логика консервативна:
        не копируем только явно пустые плейсхолдеры/мусор.
        """
        try:
            if getattr(shape, "is_placeholder", False):
                if getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip():
                    return True
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    return True
                if getattr(shape, "has_table", False):
                    return True
                return False

            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP):
                return True
            if getattr(shape, "has_table", False):
                return True
            if getattr(shape, "has_chart", False):
                return True
            if getattr(shape, "has_text_frame", False) and (shape.text_frame.text or "").strip():
                return True

            if shape.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.LINE,
                MSO_SHAPE_TYPE.CONNECTOR,
            ):
                return True

            return self._shape_has_visible_fill_or_line(shape)
        except Exception:
            return True
    
    def _remove_slides_after(self, prs, keep_count: int):
        """Удаляет все слайды начиная с keep_count (оставляет первые keep_count)."""
        sldIdLst = prs.slides._sldIdLst
        while len(sldIdLst) > keep_count:
            sldIdLst.remove(sldIdLst[-1])

    def _safe_int(self, v, default=0):
        try:
            if v is None:
                return default
            return int(v)
        except Exception:
            return default

    def _fill_title_slide_texts(self, slide, top_text: str, bottom_text: str):
        """Заполняет два основных текстовых блока на титульнике, сохраняя формат."""
        text_shapes = []
        for sh in slide.shapes:
            try:
                if sh.has_text_frame:
                    text_shapes.append(sh)
            except Exception:
                continue

        text_shapes.sort(key=lambda s: (self._safe_int(getattr(s, "top", 0)), self._safe_int(getattr(s, "left", 0))))

        def set_text_preserve(shape, new_text: str):
            try:
                tf = shape.text_frame
                if not tf.paragraphs:
                    return
                p0 = tf.paragraphs[0]
                runs = p0.runs
                if runs:
                    runs[0].text = new_text
                    for r in runs[1:]:
                        r.text = ""
                else:
                    p0.text = new_text
                for p in tf.paragraphs[1:]:
                    p.text = ""
            except Exception:
                pass

        if len(text_shapes) >= 1:
            set_text_preserve(text_shapes[0], top_text or "")
        if len(text_shapes) >= 2:
            set_text_preserve(text_shapes[1], bottom_text or "")

    def _extract_title_texts(self, src_slide):
        texts = []
        for sh in src_slide.shapes:
            try:
                if sh.has_text_frame:
                    t = (sh.text_frame.text or "").strip()
                    if t:
                        texts.append((self._safe_int(getattr(sh, "top", 0)), self._safe_int(getattr(sh, "left", 0)), t))
            except Exception:
                continue
        texts.sort(key=lambda x: (x[0], x[1]))
        top = texts[0][2] if len(texts) >= 1 else ""
        bottom = texts[1][2] if len(texts) >= 2 else ""
        return top, bottom

    def _clean_template_placeholders(self, slide):
        to_delete = []
        for sh in slide.shapes:
            try:
                if getattr(sh, "is_placeholder", False):
                    to_delete.append(sh)
                    continue
                if sh.has_text_frame:
                    t = (sh.text_frame.text or "").strip().lower()
                    if t in ("заголовок слайда", "подзаголовок", "текст", "title", "subtitle"):
                        to_delete.append(sh)
            except Exception:
                continue

        for sh in to_delete:
            try:
                slide.shapes._spTree.remove(sh._element)
            except Exception:
                pass

    def fix_presentation(self, out_path: str) -> str:
        """Стабильная генерация:
        - титульники 1:1 из template (слайды 1-3) + подстановка текста
        - остальные: фон/оформление из template (слайд 4) + перенос фигур из исходника
        """
        src_prs = Presentation(self.pptx_path)
        tpl_prs = Presentation(self.template_path)

        if len(tpl_prs.slides) < 4:
            raise ValueError("В template.pptx должно быть минимум 4 слайда.")

        dst_prs = Presentation(self.template_path)
        self._remove_slides_after(dst_prs, 3)

        # титульники
        title_top, title_bottom = ("", "")
        if len(src_prs.slides) > 0:
            title_top, title_bottom = self._extract_title_texts(src_prs.slides[0])

        for i in range(min(3, len(dst_prs.slides))):
            self._fill_title_slide_texts(dst_prs.slides[i], title_top, title_bottom)

        base_tpl_content = tpl_prs.slides[3]
        base_layout = base_tpl_content.slide_layout

        for src_idx in range(1, len(src_prs.slides)):
            src_slide = src_prs.slides[src_idx]
            dst_slide = dst_prs.slides.add_slide(base_layout)

            # фон/оформление из шаблона
            self.copy_slide_shapes(base_tpl_content, dst_slide)
            self._clean_template_placeholders(dst_slide)

            # контент из исходника
            self.copy_slide_shapes(src_slide, dst_slide)

        dst_prs.save(out_path)
        return out_path
    
    def _extract_title_texts(self, slide):
        """Берём текстовые блоки с первого слайда исходника, сортируем сверху-вниз."""
        blocks = []
        for sh in slide.shapes:
            if hasattr(sh, "text_frame") and sh.has_text_frame:
                txt = (sh.text_frame.text or "").strip()
                if txt:
                    blocks.append((sh.top, sh.left, txt))
        blocks.sort(key=lambda x: (x[0], x[1]))
        return [b[2] for b in blocks]

    def _fill_template_title_slide(self, tpl_slide, src_texts: list[str]):
        """Заполняем титульный слайд шаблона, не трогая оформление (шрифты/фон/размеры)."""
        if not src_texts:
            return

        tpl_text_shapes = []
        for sh in tpl_slide.shapes:
            if hasattr(sh, "text_frame") and sh.has_text_frame:
                tpl_text_shapes.append(sh)

        tpl_text_shapes.sort(key=lambda s: (s.top, s.left))

        for idx, sh in enumerate(tpl_text_shapes):
            if idx >= len(src_texts):
                break

            if idx == len(tpl_text_shapes) - 1 and len(src_texts) > len(tpl_text_shapes):
                value = "\n".join(src_texts[idx:])
            else:
                value = src_texts[idx]

            self._replace_text_keep_format(sh, value)

    def _replace_text_keep_format(self, shape, new_text: str, force_font: dict | None = None):
        """Заменяем текст, сохраняя форматирование первого run в каждом абзаце."""
        if not (hasattr(shape, "text_frame") and shape.has_text_frame):
            return

        tf = shape.text_frame
        lines = (new_text or "").splitlines() or [""]

        if len(tf.paragraphs) == 0:
            tf.text = ""

        while len(tf.paragraphs) < len(lines):
            tf.add_paragraph()
        while len(tf.paragraphs) > len(lines):
            p = tf.paragraphs[-1]._p
            p.getparent().remove(p)

        for i, line in enumerate(lines):
            p = tf.paragraphs[i]

            if force_font:
                try:
                    p.font.name = force_font.get("name")
                    p.font.size = force_font.get("size")
                    if force_font.get("bold") is not None:
                        p.font.bold = force_font.get("bold")
                except Exception:
                    pass

            if p.runs:
                p.runs[0].text = line
                if force_font:
                    try:
                        p.runs[0].font.name = force_font.get("name")
                        p.runs[0].font.size = force_font.get("size")
                        if force_font.get("bold") is not None:
                            p.runs[0].font.bold = force_font.get("bold")
                    except Exception:
                        pass
                for r in p.runs[1:]:
                    r.text = ""
            else:
                r = p.add_run()
                r.text = line
                if force_font:
                    try:
                        r.font.name = force_font.get("name")
                        r.font.size = force_font.get("size")
                        if force_font.get("bold") is not None:
                            r.font.bold = force_font.get("bold")
                    except Exception:
                        pass

    # -----------------------------
    # Content slides
    # -----------------------------
    def _prepare_destination_content_slide(self, dst_slide):
        """Оставляем фон/оформление макета, но убираем placeholders (чтобы не было лишних фигур)."""
        to_remove = []
        for sh in dst_slide.shapes:
            try:
                if getattr(sh, "is_placeholder", False):
                    to_remove.append(sh)
            except Exception:
                pass

        for sh in reversed(to_remove):
            try:
                el = sh.element
                el.getparent().remove(el)
            except Exception:
                pass

    def copy_slide_shapes(self, src_slide, dst_slide):
        """Копируем ВСЕ фигуры со слайда исходника на целевой слайд (сохраняя размеры/позиции/таблицы/картинки)."""
        for src_shape in src_slide.shapes:
            try:
                # Skip invisible / "junk" shapes to avoid bloating and visual garbage
                if not self._is_shape_meaningful(src_shape):
                    continue
                self.copy_shape(src_shape, dst_slide)
            except Exception as e:
                print(f"    ⚠️  Ошибка копирования shape: {e}")

    def copy_shape(self, src_shape, dst_slide):
        # 1) Таблица
        if getattr(src_shape, "has_table", False):
            self._copy_table(src_shape, dst_slide)
            return

        # 2) Картинка (PIC) — копируем XML + rels, чтобы сохранить crop/rotate/effects
        if src_shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(src_shape, "image"):
            self._copy_picture_xml(src_shape, dst_slide)
            return

        # 3) Текстовое поле / placeholder с текстом
        if hasattr(src_shape, "text_frame") and getattr(src_shape, "has_text_frame", False):
            txt = (src_shape.text_frame.text or "").strip()
            if txt:
                self._copy_textbox_keep_size(src_shape, dst_slide)
            return

        # 4) Автофигуры/прочее — если есть заливка/линия, лучше переносить как XML (без rels)
        if src_shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            self._copy_shape_element(src_shape, dst_slide)
            return

        # 5) Группы — рекурсивно
        if src_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in src_shape.shapes:
                try:
                    if not self._is_shape_meaningful(sub):
                        continue
                except Exception:
                    pass
                self.copy_shape(sub, dst_slide)
            return

    # -----------------------------
    # Text copy (keep font sizes from source, change name to Montserrat)
    # -----------------------------
    def _copy_textbox_keep_size(self, src_shape, dst_slide):
        tb = dst_slide.shapes.add_textbox(src_shape.left, src_shape.top, src_shape.width, src_shape.height)
        tb.text_frame.word_wrap = True
        tb.text_frame.clear()

        src_tf = src_shape.text_frame
        for p_idx, src_p in enumerate(src_tf.paragraphs):
            dst_p = tb.text_frame.paragraphs[0] if p_idx == 0 else tb.text_frame.add_paragraph()
            # переносим базовые свойства абзаца
            try:
                dst_p.level = src_p.level
            except Exception:
                pass
            try:
                dst_p.alignment = src_p.alignment
            except Exception:
                pass

            # runs
            if not src_p.runs:
                dst_p.text = src_p.text
                # применим Montserrat к пустому run, если появился
                for r in dst_p.runs:
                    self._apply_run_font_from_source(r, None)
                continue

            # чистим автосозданный пустой run
            dst_p.text = ""
            for src_r in src_p.runs:
                dst_r = dst_p.add_run()
                dst_r.text = src_r.text

                # copy bold/italic/underline/color + SIZE from source
                self._apply_run_font_from_source(dst_r, src_r)

    def _apply_run_font_from_source(self, dst_run, src_run, src_paragraph=None):
        """Переносит стиль текста, сохраняя размер как в исходнике.

        Часто размер задаётся на уровне paragraph.font.size, а у runs бывает None.
        Если не подхватить это — берётся дефолт из шаблона (например, 18 вместо 14).
        """
        try:
            dst_run.font.bold = src_run.font.bold
            dst_run.font.italic = src_run.font.italic
            dst_run.font.underline = src_run.font.underline

            if src_run.font.color and src_run.font.color.rgb:
                dst_run.font.color.rgb = src_run.font.color.rgb

            size = src_run.font.size
            if size is None and src_paragraph is not None:
                try:
                    size = src_paragraph.font.size
                except Exception:
                    size = None
            dst_run.font.size = size

            dst_run.font.name = "Montserrat"
        except Exception:
            pass
            
    def _copy_table(self, src_shape, dst_slide):
        rows = len(src_shape.table.rows)
        cols = len(src_shape.table.columns)

        left, top, width, height = src_shape.left, src_shape.top, src_shape.width, src_shape.height
        graphic_frame = dst_slide.shapes.add_table(rows, cols, left, top, width, height)
        dst_table = graphic_frame.table
        src_table = src_shape.table

        # ширины колонок
        try:
            for c in range(cols):
                dst_table.columns[c].width = src_table.columns[c].width
        except Exception:
            pass
        # высоты строк
        try:
            for r in range(rows):
                dst_table.rows[r].height = src_table.rows[r].height
        except Exception:
            pass

        for r in range(rows):
            for c in range(cols):
                src_cell = src_table.cell(r, c)
                dst_cell = dst_table.cell(r, c)

                # текст
                dst_tf = dst_cell.text_frame
                dst_tf.clear()
                src_tf = src_cell.text_frame

                for p_idx, src_p in enumerate(src_tf.paragraphs):
                    dst_p = dst_tf.paragraphs[0] if p_idx == 0 else dst_tf.add_paragraph()
                    try:
                        dst_p.alignment = src_p.alignment
                    except Exception:
                        pass
                    dst_p.text = ""

                    if not src_p.runs:
                        dst_p.text = src_p.text
                        continue

                    for src_r in src_p.runs:
                        dst_r = dst_p.add_run()
                        dst_r.text = src_r.text
                        # шрифт Montserrat, но размер и жирность из исходника
                        dst_r.font.name = "Montserrat"
                        try:
                            if src_r.font.size is not None:
                                dst_r.font.size = src_r.font.size
                        except Exception:
                            pass
                        for attr in ("bold", "italic", "underline"):
                            try:
                                if getattr(src_r.font, attr) is not None:
                                    setattr(dst_r.font, attr, getattr(src_r.font, attr))
                            except Exception:
                                pass
                        try:
                            if src_r.font.color and src_r.font.color.rgb is not None:
                                dst_r.font.color.rgb = src_r.font.color.rgb
                        except Exception:
                            pass

    # -----------------------------
    # Picture copy via XML + rels (preserves crop/rotate/effects)
    # -----------------------------
    def _copy_picture_xml(self, src_pic_shape, dst_slide):
        from copy import deepcopy
        import io

        # 1) добавляем/находим image part в dst слайде
        image_blob = src_pic_shape.image.blob
        image_part, rId = dst_slide.part.get_or_add_image_part(io.BytesIO(image_blob))

        # 2) копируем XML pic и подменяем rId
        pic = deepcopy(src_pic_shape._element)
        # найти blip
        blip = pic.xpath('.//a:blip')[0]
        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', rId)

        # 3) вставляем в spTree перед extLst (если есть), иначе в конец
        spTree = dst_slide.shapes._spTree
        extLst = spTree.xpath('./p:extLst')
        if extLst:
            spTree.insert(spTree.index(extLst[0]), pic)
        else:
            spTree.append(pic)

    # -----------------------------
    # Generic shape element copy (for simple autoshapes)
    # -----------------------------
    def _copy_shape_element(self, src_shape, dst_slide):
        from copy import deepcopy
        el = deepcopy(src_shape._element)
        spTree = dst_slide.shapes._spTree
        extLst = spTree.xpath('./p:extLst')
        if extLst:
            spTree.insert(spTree.index(extLst[0]), el)
        else:
            spTree.append(el)

    # -----------------------------
    # Utilities: keep only first N slides in template
    # -----------------------------
    def _keep_only_first_n_slides(self, prs, n: int):
        # Удаляем слайды с конца, чтобы индексы не скакали
        while len(prs.slides) > n:
            self._delete_slide(prs, len(prs.slides) - 1)

    def _delete_slide(self, prs, index: int):
        # index — 0-based
        sldIdLst = prs.slides._sldIdLst
        sldId = sldIdLst[index]
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)